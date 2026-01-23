"""PPT解析服务"""
import subprocess
import asyncio
from pathlib import Path
from typing import List, Dict, Any
from pptx import Presentation
from app.config import settings


def extract_text_from_ppt(file_path: str) -> List[Dict[str, Any]]:
    """
    从 .pptx 文件中提取每页的文本内容

    Args:
        file_path: PPT 文件路径

    Returns:
        每页文本内容的列表，每项包含 page_num 和 content
    """
    path = Path(file_path)

    if not path.exists():
        raise FileNotFoundError(f"文件不存在: {file_path}")

    if not path.suffix.lower() == '.pptx':
        raise ValueError("只支持 .pptx 格式的 PowerPoint 文件")

    prs = Presentation(file_path)
    slides = []

    for idx, slide in enumerate(prs.slides, start=1):
        slide_text = []
        slide_images = []

        for shape in slide.shapes:
            # 提取文本框内容
            if hasattr(shape, "text_frame") and shape.text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_text.append(text)

            # 检查是否有图片
            if hasattr(shape, "shape_type"):
                # 检查是否是图片类型
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    if hasattr(shape, "name"):
                        slide_images.append(shape.name)

        slides.append({
            "page_num": idx,
            "content": "\n".join(slide_text) if slide_text else "[无文本内容]",
            "image_count": len(slide_images)
        })

    return slides


def extract_all_text_from_ppt(file_path: str) -> str:
    """
    提取所有文本内容（兼容旧接口）

    Args:
        file_path: PPT 文件路径

    Returns:
        格式化的文本字符串，每页幻灯片用分隔符区分
    """
    slides = extract_text_from_ppt(file_path)
    result = []
    for slide in slides:
        result.append(f"=== 第 {slide['page_num']} 页 ===\n{slide['content']}")
    return "\n\n".join(result)


def get_ppt_info(file_path: str) -> Dict[str, Any]:
    """
    获取 PPT 文件的基本信息

    Args:
        file_path: PPT 文件路径

    Returns:
        包含文件信息的字典
    """
    path = Path(file_path)

    if not path.exists():
        raise FileNotFoundError(f"文件不存在: {file_path}")

    prs = Presentation(file_path)

    return {
        "filename": path.name,
        "slide_count": len(prs.slides),
        "file_size": path.stat().st_size,
        "slides": extract_text_from_ppt(file_path)
    }


def generate_slide_screenshots(ppt_path: str, task_id: str) -> List[Dict[str, Any]]:
    """
    使用 LibreOffice 将 PPT 转换为图片截图

    Args:
        ppt_path: PPT 文件路径
        task_id: 任务 ID（用于生成截图文件夹名）

    Returns:
        包含每页截图路径的列表
    """
    path = Path(ppt_path)
    if not path.exists():
        raise FileNotFoundError(f"文件不存在: {ppt_path}")

    # 确保截图目录存在
    screenshot_dir = settings.static_dir / "screenshots" / task_id
    screenshot_dir.mkdir(parents=True, exist_ok=True)

    # 使用 LibreOffice 将 PPT 转换为 PDF，然后再转为图片
    # 先转换为 PDF
    pdf_path = screenshot_dir / f"{path.stem}.pdf"

    # LibreOffice 命令
    cmd = [
        "soffice",
        "--headless",
        "--convert-to", "pdf",
        "--outdir", str(screenshot_dir),
        str(ppt_path)
    ]

    try:
        result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
        if result.returncode != 0:
            raise RuntimeError(f"LibreOffice 转换失败: {result.stderr}")

        # 如果 PDF 不存在，检查是否有其他问题
        if not pdf_path.exists():
            raise RuntimeError(f"PDF 文件未生成: {pdf_path}")

    except subprocess.TimeoutExpired:
        raise RuntimeError("LibreOffice 转换超时")

    # 使用 pdftoppm 将 PDF 转为图片
    images = []
    if pdf_path.exists():
        try:
            # pdftoppm -png input.pdf output_prefix
            # 不使用 -singlefile，这样会生成所有页面，文件名格式为 prefix-1.png, prefix-2.png 等
            result = subprocess.run(
                ["pdftoppm", "-png", str(pdf_path), str(screenshot_dir / path.stem)],
                capture_output=True,
                timeout=120
            )

            if result.returncode == 0:
                # 查找生成的图片文件
                for png_file in sorted(screenshot_dir.glob(f"{path.stem}-*.png")):
                    rel_path = f"/static/screenshots/{task_id}/{png_file.name}"
                    # 从文件名提取页码 (filename-1.png)
                    page_num = int(png_file.stem.split('-')[-1])
                    images.append({
                        "page_num": page_num,
                        "screenshot_path": rel_path
                    })
            else:
                raise RuntimeError(f"pdftoppm 失败: {result.stderr.decode()}")

        except FileNotFoundError:
            raise RuntimeError("未找到 pdftoppm，请安装 poppler-utils")

    # 清理临时 PDF 文件
    if pdf_path.exists():
        pdf_path.unlink()

    return images


async def generate_slide_screenshots_async(ppt_path: str, task_id: str) -> List[Dict[str, Any]]:
    """
    异步生成幻灯片截图
    """
    loop = asyncio.get_event_loop()
    return await loop.run_in_executor(
        None,
        lambda: generate_slide_screenshots(ppt_path, task_id)
    )
